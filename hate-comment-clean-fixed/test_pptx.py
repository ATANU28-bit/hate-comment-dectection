from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

try:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    print("Testing ROUNDED_RECTANGLE .solid()")
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 100, 100, 100, 100)
    
    try:
        shape.fill.solid()
        print("Success .solid()")
    except Exception as e:
        print(f"Failed .solid(): {e}")

except Exception as e:
    print(f"Global error: {e}")
